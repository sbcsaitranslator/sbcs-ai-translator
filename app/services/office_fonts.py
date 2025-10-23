# from __future__ import annotations
# from typing import Iterable
# from io import BytesIO

# def _font_for_lang(target_lang: str) -> str:
#     t = (target_lang or "").lower()
#     if t.startswith(("ja", "jp", "ko", "zh", "zh-hans", "zh-hant")):
#         return "Meiryo UI"
#     return "Calibri"

# def set_docx_font(data: bytes, font_name: str) -> bytes:
#     try:
#         from docx import Document
#         from docx.oxml.ns import qn
#     except Exception:
#         return data
#     doc = Document(BytesIO(data))

#     def _apply_run(run):
#         try:
#             run.font.name = font_name
#             r = run._element.rPr
#             if r is not None:
#                 if r.rFonts is None:
#                     r._new_rPr()
#                 r.rFonts.set(qn('w:ascii'), font_name)
#                 r.rFonts.set(qn('w:hAnsi'), font_name)
#                 r.rFonts.set(qn('w:cs'), font_name)
#                 r.rFonts.set(qn('w:eastAsia'), font_name)
#         except Exception:
#             pass

#     for p in doc.paragraphs:
#         for r in p.runs:
#             _apply_run(r)
#     for table in doc.tables:
#         for row in table.rows:
#             for cell in row.cells:
#                 for p in cell.paragraphs:
#                     for r in p.runs:
#                         _apply_run(r)

#     try:
#         doc.styles['Normal'].font.name = font_name
#     except Exception:
#         pass

#     out = BytesIO(); doc.save(out); return out.getvalue()

# def set_pptx_font(data: bytes, font_name: str) -> bytes:
#     try:
#         from pptx import Presentation
#         from pptx.enum.shapes import MSO_SHAPE_TYPE
#         from pptx.oxml.xmlchemy import OxmlElement
#         from pptx.oxml.ns import qn
#     except Exception:
#         return data

#     def _force_run(run):
#         try:
#             rPr = run._r.get_or_add_rPr()
#             latin = rPr.find(qn('a:latin')) or OxmlElement('a:latin'); rPr.append(latin) if latin.getparent() is None else None
#             ea    = rPr.find(qn('a:ea'))    or OxmlElement('a:ea');    rPr.append(ea)    if ea.getparent()    is None else None
#             cs    = rPr.find(qn('a:cs'))    or OxmlElement('a:cs');    rPr.append(cs)    if cs.getparent()    is None else None
#             latin.set('typeface', font_name)
#             ea.set('typeface', font_name)
#             cs.set('typeface', font_name)
#             run.font.name = font_name
#         except Exception:
#             pass

#     def _apply_shape(shape):
#         try:
#             if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
#                 for p in shape.text_frame.paragraphs:
#                     for r in p.runs:
#                         _force_run(r)
#             if getattr(shape, "has_table", False) and shape.has_table:
#                 for row in shape.table.rows:
#                     for cell in row.cells:
#                         if hasattr(cell, "text_frame") and cell.text_frame:
#                             for p in cell.text_frame.paragraphs:
#                                 for r in p.runs:
#                                     _force_run(r)
#             if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
#                 for s in shape.shapes:
#                     _apply_shape(s)
#         except Exception:
#             pass

#     prs = Presentation(BytesIO(data))
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             _apply_shape(shape)

#     # coba set juga default text styles
#     try:
#         ts = prs.slide_master.text_styles
#         if ts:
#             for cat in (ts.title, ts.body, ts.other):
#                 if not cat:
#                     continue
#                 for i in range(1, 10):
#                     f = getattr(cat, f'level{i}_font', None)
#                     if f:
#                         f.name = font_name
#     except Exception:
#         pass

#     out = BytesIO(); prs.save(out); return out.getvalue()

# def enforce_fonts_by_lang(name_or_ext: str, bin_data: bytes, target_lang: str) -> bytes:
#     font = _font_for_lang(target_lang)
#     e = (name_or_ext or "").lower()
#     e = e if "." in e else f".{e}"  # biar "pptx" juga match
#     try:
#         if e.endswith(".docx") or e.endswith(".doc"):
#             return set_docx_font(bin_data, font)
#         if e.endswith(".pptx") or e.endswith(".ppt"):
#             return set_pptx_font(bin_data, font)
#         return bin_data
#     except Exception:
#         return bin_data


# from __future__ import annotations
# from io import BytesIO

# def _font_for_lang(target_lang: str) -> str:
#     """
#     Pilih font default yang aman untuk target bahasa.
#     - CJK → Meiryo UI (tersedia luas & mendukung East Asia)
#     - Lainnya → Calibri
#     """
#     t = (target_lang or "").lower()
#     if t.startswith(("ja", "jp", "ko", "zh", "zh-hans", "zh-hant")):
#         return "Meiryo UI"
#     return "Calibri"


# def set_docx_font(data: bytes, font_name: str) -> bytes:
#     """
#     Ubah font pada DOCX secara aman:
#     - Set run.font.name
#     - Pastikan w:rPr/w:rFonts ada lalu isi ascii/hAnsi/cs/eastAsia
#     """
#     try:
#         from docx import Document
#         from docx.oxml.xmlchemy import OxmlElement
#         from docx.oxml.ns import qn
#     except Exception:
#         # python-docx tidak tersedia → kembalikan apa adanya
#         return data

#     try:
#         doc = Document(BytesIO(data))
#     except Exception:
#         return data

#     def _apply_run(run) -> None:
#         try:
#             # API resmi
#             run.font.name = font_name

#             # Pastikan rPr & rFonts ada lalu set mapping font
#             r = run._element  # <w:r>
#             rPr = r.rPr
#             if rPr is None:
#                 rPr = OxmlElement("w:rPr")
#                 r.insert(0, rPr)

#             rFonts = rPr.rFonts
#             if rFonts is None:
#                 rFonts = OxmlElement("w:rFonts")
#                 rPr.append(rFonts)

#             rFonts.set(qn("w:ascii"), font_name)
#             rFonts.set(qn("w:hAnsi"), font_name)
#             rFonts.set(qn("w:cs"), font_name)
#             rFonts.set(qn("w:eastAsia"), font_name)
#         except Exception:
#             # Jangan sampai gagal total hanya karena 1 run
#             pass

#     # Paragraph runs
#     for p in doc.paragraphs:
#         for r in p.runs:
#             _apply_run(r)

#     # Table runs
#     for table in doc.tables:
#         for row in table.rows:
#             for cell in row.cells:
#                 for p in cell.paragraphs:
#                     for r in p.runs:
#                         _apply_run(r)

#     # Default style (kalau ada)
#     try:
#         doc.styles["Normal"].font.name = font_name
#     except Exception:
#         pass

#     out = BytesIO()
#     doc.save(out)
#     return out.getvalue()


# def set_pptx_font(data: bytes, font_name: str) -> bytes:
#     """
#     Ubah font pada PPTX secara aman:
#     - **Tanpa** menyuntik OXML mentah (a:latin/a:ea/a:cs) → ini yang sering memicu 'Repair'
#     - Gunakan hanya API python-pptx:
#       * run.font.name
#       * recusive ke group shapes
#       * tables → cell.text_frame
#       * set default via slide_master.text_styles bila tersedia
#     """
#     try:
#         from pptx import Presentation
#         from pptx.enum.shapes import MSO_SHAPE_TYPE
#     except Exception:
#         return data

#     try:
#         prs = Presentation(BytesIO(data))
#     except Exception:
#         # File bukan PPTX valid atau rusak sejak awal → jangan diutak-atik
#         return data

#     def _apply_text_frame(tf) -> None:
#         if not tf:
#             return
#         for p in tf.paragraphs:
#             # Set default paragraph font (jika ada)
#             try:
#                 if hasattr(p, "font") and p.font:
#                     p.font.name = font_name
#             except Exception:
#                 pass
#             # Set run level
#             for r in p.runs:
#                 try:
#                     r.font.name = font_name
#                 except Exception:
#                     pass

#     def _walk_shape(shape) -> None:
#         try:
#             # Text box / placeholder
#             if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
#                 _apply_text_frame(shape.text_frame)

#             # Table
#             if getattr(shape, "has_table", False) and shape.has_table:
#                 for row in shape.table.rows:
#                     for cell in row.cells:
#                         if getattr(cell, "text_frame", None):
#                             _apply_text_frame(cell.text_frame)

#             # Group: rekursif
#             if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
#                 for s in shape.shapes:
#                     _walk_shape(s)

#             # Catatan: chart/smartart kadang punya teks, tapi python-pptx belum expose
#             # API stabil untuk semua kasus. Kita hindari manipulasi XML mentah demi
#             # menjaga integritas file.
#         except Exception:
#             pass

#     # Semua slide
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             _walk_shape(shape)

#     # Default via master text styles (jika ada)
#     try:
#         ts = getattr(prs.slide_master, "text_styles", None)
#         if ts:
#             for cat in (ts.title, ts.body, ts.other):
#                 if not cat:
#                     continue
#                 for lvl in range(1, 10):
#                     f = getattr(cat, f"level{lvl}_font", None)
#                     if f:
#                         try:
#                             f.name = font_name
#                         except Exception:
#                             pass
#     except Exception:
#         pass

#     out = BytesIO()
#     try:
#         prs.save(out)
#     except Exception:
#         # Jika terjadi error saat save, kembalikan data asli (fail-safe)
#         return data
#     return out.getvalue()


# def enforce_fonts_by_lang(name_or_ext: str, bin_data: bytes, target_lang: str) -> bytes:
#     """
#     Dispatcher sederhana berbasis ekstensi. Pertahankan nama & perilaku.
#     """
#     font = _font_for_lang(target_lang)
#     e = (name_or_ext or "").lower()
#     e = e if "." in e else f".{e}"  # biar "pptx" juga match
#     try:
#         if e.endswith(".docx") or e.endswith(".doc"):
#             return set_docx_font(bin_data, font)
#         if e.endswith(".pptx") or e.endswith(".ppt"):
#             return set_pptx_font(bin_data, font)
#         return bin_data
#     except Exception:
#         return bin_data


from __future__ import annotations
from io import BytesIO

# =========================
# 1) Language → Font mapper
# =========================
def _font_for_lang(target_lang: str) -> str:
    """
    CJK → Meiryo UI, lainnya → Calibri
    """
    t = (target_lang or "").lower()
    if t.startswith(("ja", "jp", "ko", "zh", "zh-hans", "zh-hant")):
        return "Meiryo UI"
    return "Calibri"


# =========================
# 2) DOCX font enforcer
# =========================
def set_docx_font(data: bytes, font_name: str) -> bytes:
    try:
        from docx import Document
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
    except Exception:
        return data

    try:
        doc = Document(BytesIO(data))
    except Exception:
        return data

    def _ensure_rPr(elem):
        rPr = elem.rPr
        if rPr is None:
            rPr = OxmlElement('w:rPr')
            elem._element.append(rPr)
        return rPr

    def _ensure_rFonts(rPr):
        rFonts = rPr.find(qn('w:rFonts'))
        if rFonts is None:
            rFonts = OxmlElement('w:rFonts')
            rPr.append(rFonts)
        # bersihkan theme agar tidak override
        for a in ('asciiTheme', 'hAnsiTheme', 'eastAsiaTheme', 'cstheme'):
            k = qn(f'w:{a}')
            if rFonts.get(k):
                del rFonts.attrib[k]
        return rFonts

    def _force_run_font(run):
        try:
            run.font.name = font_name
            rPr = run._element.rPr
            if rPr is None:
                rPr = OxmlElement('w:rPr')
                run._element.append(rPr)
            rFonts = _ensure_rFonts(rPr)
            rFonts.set(qn('w:ascii'), font_name)
            rFonts.set(qn('w:hAnsi'), font_name)
            rFonts.set(qn('w:eastAsia'), font_name)  # PENTING utk CJK
            rFonts.set(qn('w:cs'), font_name)
        except Exception:
            pass

    def _apply_paragraph(p):
        # style paragraph (default level)
        try:
            if getattr(p, "style", None) and getattr(p.style, "font", None):
                p.style.font.name = font_name
        except Exception:
            pass
        for r in p.runs:
            _force_run_font(r)

    # Seluruh paragraf top-level
    for p in doc.paragraphs:
        _apply_paragraph(p)

    # Di dalam tabel
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _apply_paragraph(p)

    # Defaults dokumen (docDefaults) agar placeholder/teks baru ikut
    try:
        styles_elm = doc.styles.element
        docDefaults = styles_elm.find(qn('w:docDefaults'))
        if docDefaults is None:
            docDefaults = OxmlElement('w:docDefaults')
            styles_elm.append(docDefaults)

        rPrDefault = docDefaults.find(qn('w:rPrDefault'))
        if rPrDefault is None:
            rPrDefault = OxmlElement('w:rPrDefault')
            docDefaults.append(rPrDefault)

        rPr = rPrDefault.find(qn('w:rPr'))
        if rPr is None:
            rPr = OxmlElement('w:rPr')
            rPrDefault.append(rPr)

        rFonts = _ensure_rFonts(rPr)
        rFonts.set(qn('w:ascii'), font_name)
        rFonts.set(qn('w:hAnsi'), font_name)
        rFonts.set(qn('w:eastAsia'), font_name)
        rFonts.set(qn('w:cs'), font_name)

        # Style "Normal" juga diset, sebagai jaring pengaman
        try:
            normal = doc.styles['Normal']
            normal.font.name = font_name
            # pastikan ada rFonts di style Normal
            n_rPr = normal.element.get_or_add_rPr()
            n_rFonts = n_rPr.find(qn('w:rFonts'))
            if n_rFonts is None:
                n_rFonts = OxmlElement('w:rFonts')
                n_rPr.append(n_rFonts)
            n_rFonts.set(qn('w:ascii'), font_name)
            n_rFonts.set(qn('w:hAnsi'), font_name)
            n_rFonts.set(qn('w:eastAsia'), font_name)
            n_rFonts.set(qn('w:cs'), font_name)
        except Exception:
            pass
    except Exception:
        pass

    out = BytesIO()
    try:
        doc.save(out)
        return out.getvalue()
    except Exception:
        return data


# =========================
# 3) PPTX font enforcer
# =========================
def set_pptx_font(data: bytes, font_name: str) -> bytes:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
        from pptx.oxml.xmlchemy import OxmlElement
    except Exception:
        return data

    try:
        prs = Presentation(BytesIO(data))
    except Exception:
        return data

    def _ensure_font_nodes(rPr):
        latin = rPr.find(qn('a:latin'))
        if latin is None:
            latin = OxmlElement('a:latin'); rPr.append(latin)
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = OxmlElement('a:ea'); rPr.append(ea)
        cs = rPr.find(qn('a:cs'))
        if cs is None:
            cs = OxmlElement('a:cs'); rPr.append(cs)
        return latin, ea, cs

    def _force_run_font(run):
        try:
            run.font.name = font_name
            rPr = run._r.get_or_add_rPr()
            latin, ea, cs = _ensure_font_nodes(rPr)
            latin.set('typeface', font_name)
            ea.set('typeface', font_name)   # PENTING utk CJK
            cs.set('typeface', font_name)
        except Exception:
            pass

    def _apply_text_frame(tf):
        if not tf:
            return
        for p in tf.paragraphs:
            try:
                if getattr(p, "font", None):
                    p.font.name = font_name
            except Exception:
                pass
            for r in p.runs:
                _force_run_font(r)

    def _walk_shape(shape):
        try:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                _apply_text_frame(shape.text_frame)

            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if getattr(cell, "text_frame", None):
                            _apply_text_frame(cell.text_frame)

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
                for s in shape.shapes:
                    _walk_shape(s)
        except Exception:
            pass

    for slide in prs.slides:
        for shape in slide.shapes:
            _walk_shape(shape)

    # Master text styles → title/body/other
    try:
        ts = getattr(prs.slide_master, "text_styles", None)
        if ts:
            for cat in (ts.title, ts.body, ts.other):
                if not cat:
                    continue
                for lvl in range(1, 10):
                    f = getattr(cat, f"level{lvl}_font", None)
                    if f:
                        try:
                            f.name = font_name
                        except Exception:
                            pass
    except Exception:
        pass

    out = BytesIO()
    try:
        prs.save(out)
        return out.getvalue()
    except Exception:
        return data


# =========================
# 4) Dispatcher by extension
# =========================
def enforce_fonts_by_lang(name_or_ext: str, bin_data: bytes, target_lang: str) -> bytes:
    """
    JA/KO/ZH → Meiryo UI, lainnya → Calibri.
    Otomatis pilih handler berdasar ekstensi.
    """
    font = _font_for_lang(target_lang)
    e = (name_or_ext or "").lower()
    e = e if "." in e else f".{e}"
    try:
        if e.endswith((".docx", ".doc")):
            return set_docx_font(bin_data, font)
        if e.endswith((".pptx", ".ppt")):
            return set_pptx_font(bin_data, font)
        return bin_data
    except Exception:
        return bin_data
